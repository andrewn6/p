from sanic import Sanic
from sanic.log import logger
from sanic import response
from sanic_cors import CORS

from pdfminer.high_level import extract_text
from pptx import Presentation

import os
import uuid
import json
import time

from nltk.tokenize import word_tokenize
import nltk
from nltk.corpus import stopwords

from io import BytesIO

import spacy
import transformers


app = Sanic("Summarizer")
CORS(app)

model_name = "facebook/bart-base"

tokenizer = transformers.AutoTokenizer.from_pretrained(model_name)
model = transformers.AutoModelForSeq2SeqLM.from_pretrained(model_name)

nlp = spacy.load("en_core_web_sm")
nltk.download('stopwords')

def get_pdf_text(file_bytes: BytesIO):
    text = extract_text(file_bytes)
    return text

def get_pptx_text(file_bytes: BytesIO):
  prs = Presentation(file_bytes)
  text = ""
  for slide in prs.slides:
    for shape in slide.shapes:
      if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
          for run in paragraph.runs:
            text += run.text

  return text

def get_file_text(supplied_file):
    match get_file_ext(supplied_file.name):
        case "pdf":
            return get_pdf_text(BytesIO(supplied_file.body))
        case "pptx":
            return get_pptx_text(BytesIO(supplied_file.body))
        case _:
            return False

def clean_text(text):
  stop_words = set(stopwords.words('english'))

  text = text.lower()

  word_tokens = word_tokenize(text)

  filtered_text = [word for word in word_tokens if word.isalnum() and not word in stop_words]

  return ' '.join(filtered_text).strip()

def summarize_text(text: str): 
    doc = nlp(text)
    sentences = [clean_text(sent.text) for sent in doc.sents]
    processed_text = "\n".join(sentences)

    max_chunk_size = 2048
    text_chunks = [processed_text[i: i + max_chunk_size] for i in range(0, len(processed_text), max_chunk_size)]
    summaries = []

    for chunk in text_chunks:
         inputs = tokenizer.encode(chunk, return_tensors="pt", truncation=True)
         summary_ids = model.generate(inputs, num_beams=4, max_length=100, min_length=10, length_penalty=2.0, early_stopping=True)
         summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
         summaries.append(summary)

    completed_summary = " ".join(summaries)

    return completed_summary

def remove_non_latin1_chars(text: str):
    return text.encode('latin-1', 'ignore').decode('latin-1')

def get_file_ext(name: str):
    return name.rsplit(".", 1)[1].lower()
def strip_file_ext(name: str):
    filename = name.rsplit(".")
    filename.pop()
    return ".".join(filename)

def write_summarization(file_path, text, file):
    if not os.path.exists('output'):
        os.makedirs('output')
    output_path = os.path.join("output", file_path)

    f = open(output_path, "w")
    f.write(json.dumps({ 
        'text': text, 
        'name': strip_file_ext(file.name), 
        'ext': get_file_ext(file.name),
        'date': round(time.time() * 1000),
    }))
    f.close()

valid_files = ["pdf", "pptx"]
def valid_file_in_request(request):
    if not request.files or 'file' not in request.files:
        return 'No file uploaded!'

    supplied_file = request.files.get('file')
    file_type = get_file_ext(supplied_file.name)

    if not supplied_file or file_type not in valid_files:
        return 'Invalid file type ' + file_type

    return supplied_file


@app.route('/summarize', methods=['POST'])
async def summarize(request):
    is_valid_file_in_request = valid_file_in_request(request)
    if type(is_valid_file_in_request) == str:
        return response.json({"message": is_valid_file_in_request}, 400)

    supplied_file = is_valid_file_in_request

    try:
        text = get_file_text(supplied_file)
        # This should have been avoided before
        if (text == False): raise Exception("Invalid file extension")
        summary = summarize_text(text)

        unique_id = str(uuid.uuid4())
        output_text_path = f"{unique_id}.json"

        write_summarization(output_text_path, summary, supplied_file)

        return response.json({"id": unique_id}, 200)

    except Exception as e:
        logger.error(f"Error during processing: {e}")
        return response.json({"message": "An error occured durring processing"}, 500)


@app.route('/summarization/<id>', methods=['GET'])
async def get_summarization(_request, id):
    try:
        path = f"output/{id}.json"
        f = open(path, "r")
        info = json.loads(f.read())
        f.close()
        return response.json({"text": info["text"], "name": info["name"], "date": info["date"], "ext": info["ext"]})

    except FileNotFoundError as e:
        return response.json({"message": "File not found. It may have been deleted"}, 404)


@app.route('/pdf-length', methods=['GET'])
async def pdf_length(request):
    is_pdf_in_request = valid_file_in_request(request)
    if type(is_pdf_in_request) == str:
        return response.json({"message": is_pdf_in_request}, 400)

    length = len(get_pdf_text(is_pdf_in_request))

    return response.json({"length": length})

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=3000)
